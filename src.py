import streamlit as st
import spacy
from transformers import pipeline, T5Tokenizer, T5ForConditionalGeneration
import matplotlib.pyplot as plt
import networkx as nx
import PyPDF2
from pptx import Presentation
from io import BytesIO
import os
from pathlib import Path
import tempfile
import re

# Streamlit page configuration
st.set_page_config(page_title="Study Notes Generator", layout="wide")

# Load NLP model
try:
    nlp = spacy.load("en_core_web_sm")
except OSError:
    st.error("Spacy model 'en_core_web_sm' not found. Please install it using: python -m spacy download en_core_web_sm")
    st.stop()

# Load transformer pipelines
try:
    summarizer = pipeline("summarization", model="facebook/bart-large-cnn")
except Exception as e:
    st.error(f"Error loading summarization model: {e}")
    st.stop()

# Load T5 model for question generation
try:
    tokenizer = T5Tokenizer.from_pretrained("valhalla/t5-small-qg-hl")
    qg_model = T5ForConditionalGeneration.from_pretrained("valhalla/t5-small-qg-hl")
except Exception as e:
    st.warning(f"Error loading question-generation model: {e}")
    qg_model, tokenizer = None, None

# Step 1: Extract text from file
def extract_text_from_file(uploaded_file):
    """Extract text from PDF or PPTX files."""
    try:
        file_extension = Path(uploaded_file.name).suffix.lower()
        file_bytes = uploaded_file.read()
        
        with tempfile.NamedTemporaryFile(delete=False, suffix=file_extension) as tmp_file:
            tmp_file.write(file_bytes)
            tmp_file_path = tmp_file.name
        
        if file_extension == '.pdf':
            with open(tmp_file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                text = ""
                for page in pdf_reader.pages:
                    extracted = page.extract_text()
                    if extracted:
                        text += extracted
                return text.strip() or None
        
        elif file_extension == '.pptx':
            presentation = Presentation(tmp_file_path)
            text = ""
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text += shape.text + "\n"
            return text.strip() or None
        
        else:
            raise ValueError("Unsupported file format. Please upload a PDF or PPTX file.")
    
    except Exception as e:
        st.error(f"Error extracting text: {e}")
        return None
    finally:
        if 'tmp_file_path' in locals():
            try:
                os.unlink(tmp_file_path)
            except Exception as e:
                st.warning(f"Failed to delete temporary file: {e}")

# Step 2: Summarize Content and Format in Points
def summarize_text(text, max_length=150, min_length=40):
    """Summarize the input text and return as bullet points."""
    try:
        if not text:
            return None
        max_chunk_size = 1000
        chunks = [text[i:i+max_chunk_size] for i in range(0, len(text), max_chunk_size)]
        summaries = []
        
        for chunk in chunks:
            summary = summarizer(chunk, max_length=max_length, min_length=min_length, do_sample=False)
            summaries.append(summary[0]['summary_text'])
        
        # Combine summaries and split into sentences for bullet points
        full_summary = " ".join(summaries)
        # Split on sentence boundaries, ensuring non-empty points
        sentences = [s.strip() for s in re.split(r'(?<=[.!?])\s+', full_summary) if s.strip()]
        return sentences
    
    except Exception as e:
        st.error(f"Error summarizing text: {e}")
        return None

# Step 3: Extract Key Points
def extract_keywords(text):
    """Extract noun phrases as key points from the text."""
    try:
        if not text:
            return []
        doc = nlp(text)
        keywords = [chunk.text for chunk in doc.noun_chunks if len(chunk.text.split()) <= 5]
        return keywords[:10]
    except Exception as e:
        st.error(f"Error extracting keywords: {e}")
        return []

# Step 4: Generate Flowchart
def create_flowchart(key_points):
    """Generate a colorful, easy-to-interpret flowchart and return as bytes."""
    try:
        if not key_points:
            st.warning("No key points provided for flowchart.")
            return None
        
        G = nx.DiGraph()
        for i in range(len(key_points) - 1):
            G.add_edge(key_points[i], key_points[i+1])
        
        fig, ax = plt.subplots(figsize=(12, 8))
        pos = nx.spring_layout(G, seed=42)  # Consistent layout
        # Colorful nodes with gradient
        node_colors = [f"C{i % 10}" for i in range(len(G.nodes))]
        # Draw nodes
        nx.draw_networkx_nodes(
            G, pos, node_color=node_colors, node_size=4000, ax=ax, edgecolors="black"
        )
        # Draw edges with distinct colors
        edge_colors = ["#FF6F61" if i % 2 == 0 else "#6B7280" for i in range(len(G.edges))]
        nx.draw_networkx_edges(G, pos, edge_color=edge_colors, width=2, ax=ax)
        # Draw labels with clear formatting
        nx.draw_networkx_labels(
            G, pos, font_size=10, font_weight="bold", font_color="white"
        )
        plt.title("Key Points Flowchart", fontsize=14, pad=20)
        plt.box(False)  # Remove box for cleaner look
        
        # Save to bytes
        buf = BytesIO()
        plt.savefig(buf, format="png", dpi=300, bbox_inches="tight", transparent=True)
        plt.close(fig)
        buf.seek(0)
        return buf
    
    except Exception as e:
        st.error(f"Error creating flowchart: {e}")
        return None

# Step 5: Question-Answer Generation
def generate_questions_and_answers(summary_sentences):
    """Generate question-answer pairs from the summary using T5 model."""
    if not qg_model or not tokenizer or not summary_sentences:
        st.warning("Question generation unavailable (model not loaded or no summary).")
        return []
    
    try:
        qa_pairs = []
        for sentence in summary_sentences[:5]:  # Limit to 5 sentences for brevity
            # Generate question
            input_text = f"generate question: {sentence}"
            inputs = tokenizer(input_text, return_tensors="pt", max_length=512, truncation=True)
            outputs = qg_model.generate(
                **inputs,
                max_length=100,
                num_beams=4,
                num_return_sequences=1,  # One question per sentence
                early_stopping=True
            )
            question = tokenizer.decode(outputs[0], skip_special_tokens=True)
            # Use sentence as answer
            answer = sentence
            qa_pairs.append({"question": question, "answer": answer})
        return qa_pairs
    except Exception as e:
        st.error(f"Error generating questions and answers: {e}")
        return []

# Streamlit UI
def main():
    st.title("NLP File Processor")
    st.write("Upload a PDF or PPTX file to summarize, extract key points, generate a flowchart, and create question-answer pairs.")
    
    # File upload
    uploaded_file = st.file_uploader("Choose a PDF or PPTX file", type=["pdf", "pptx"])
    
    if uploaded_file is not None:
        with st.spinner("Processing file..."):
            # Extract text
            text = extract_text_from_file(uploaded_file)
            if not text:
                st.error("No text extracted. Please try another file.")
                return
            
            # Summarize
            summary_sentences = summarize_text(text)
            if not summary_sentences:
                st.error("Summarization failed.")
                return
            full_summary = " ".join(summary_sentences)  # For keywords
            
            # Extract keywords
            key_points = extract_keywords(full_summary)
            
            # Create flowchart
            flowchart_buf = create_flowchart(key_points)
            
            # Generate question-answer pairs
            qa_pairs = generate_questions_and_answers(summary_sentences)
        
        # Display results
        st.subheader("Results")
        
        # Summary
        st.write("**Summary**")
        if summary_sentences:
            for sentence in summary_sentences:
                st.write(f"- {sentence}")
        else:
            st.write("No summary generated.")
        
        # Key Points
        st.write("**Key Points**")
        st.write("- " + "\n- ".join(key_points) if key_points else "None extracted.")
        
        # Flowchart
        st.write("**Flowchart**")
        if flowchart_buf:
            st.image(flowchart_buf, caption="Key Points Flowchart")
        else:
            st.write("No flowchart generated.")
        
        # Question-Answer Pairs
        st.write("**Question-Answer Pairs**")
        if qa_pairs:
            for i, qa in enumerate(qa_pairs, 1):
                st.write(f"**Q{i}.** {qa['question']}")
                st.write(f"**A{i}.** {qa['answer']}")
        else:
            st.write("No question-answer pairs generated.")

if __name__ == "__main__":
    main()

